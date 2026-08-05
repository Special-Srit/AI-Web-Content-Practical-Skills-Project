#!/usr/bin/env python3
"""냥BTI 발표 자료 빌더 — Slide Design Prompt Rev 7 준수.

내용 출처: nyangbti/docs/presentation-outline.md
실행:  <venv>/bin/python build_deck.py

Rev 7 하드 규칙:
  - 콘텐츠 상한 CONTENT_BOT = 5.08" (초과 시 assert 실패)
  - 등폭 타일 itemW = (W - (N-1)*gap) / N
  - 페이지 번호 2자리 영패딩, 표지/마지막 장 제외
  - 폰트는 Paperlogy Filled 계열 (스톡 Paperlogy는 한글 조합 미지원)
  - 이모지 금지, 그라디언트는 콘텐츠 배경 wash 하나뿐
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import pathlib

# ── 색상 토큰 (Rev 7 고정값) ────────────────────────────────────────────────
INK        = RGBColor(0x17, 0x17, 0x17)
LINK       = RGBColor(0x00, 0x70, 0xF3)
LINK_TINT  = RGBColor(0x2E, 0x86, 0xF8)
ACC_SOFT   = RGBColor(0xEB, 0xF2, 0xFF)
ACC_BORDER = RGBColor(0xC2, 0xD9, 0xFD)
CANVAS     = RGBColor(0xFF, 0xFF, 0xFF)
CANVAS_S2  = RGBColor(0xF5, 0xF5, 0xF5)
GRAD_A     = RGBColor(0xFD, 0xFD, 0xFF)
GRAD_B     = RGBColor(0xEC, 0xF1, 0xFB)
HAIRLINE   = RGBColor(0xEB, 0xEB, 0xEB)
HAIR_STR   = RGBColor(0xA1, 0xA1, 0xA1)
BODY       = RGBColor(0x4D, 0x4D, 0x4D)
MUTE       = RGBColor(0x88, 0x88, 0x88)
ON_PRI     = RGBColor(0xFF, 0xFF, 0xFF)
ON_PRI_S   = RGBColor(0xCF, 0xE2, 0xFF)

# ── 폰트 (Filled 계열 — Srit이 한글 조합 보강한 빌드) ──────────────────────
F_LIGHT  = "Paperlogy Filled Light"
F_REG    = "Paperlogy Filled Regular"
F_MED    = "Paperlogy Filled Medium"
F_SEMI   = "Paperlogy Filled SemiBold"
F_BOLD   = "Paperlogy Filled Bold"
F_XBOLD  = "Paperlogy Filled ExtraBold"
F_BLACK  = "Paperlogy Filled Black"

# ── 마스터 앵커 ────────────────────────────────────────────────────────────
SW, SH      = 10.0, 5.625
LX, W       = 0.50, 9.00
LABEL_Y     = 0.28
TITLE_Y     = 0.60
SUB_Y       = 1.15
CONTENT_Y   = 1.68
CONTENT_BOT = 5.08
FOOTER_Y    = 5.28

_violations = []


def tile(n, gap):
    """등폭 타일 규칙 — 절대 하드코딩하지 않는다."""
    return (W - (n - 1) * gap) / n


def _check(bottom, tag):
    if bottom > CONTENT_BOT + 1e-6:
        _violations.append(f"{tag}: bottom {bottom:.3f}\" > CONTENT_BOT {CONTENT_BOT}\"")


def font(run, family, size, color, spacing=None):
    """latin + ea + cs 모두 지정 — ea를 빼면 한글이 대체 폰트로 떨어진다."""
    f = run.font
    f.size = Pt(size)
    f.color.rgb = color
    f.name = family
    f.bold = False
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", family)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, family, size, color, first=False, space_before=0,
         line=None, align=None, spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    font(p.add_run() if True else None, family, size, color, spacing) if False else None
    r = p.add_run()
    r.text = text
    font(r, family, size, color, spacing)
    return p


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75, shape=MSO_SHAPE.RECTANGLE,
         radius=None):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    return sh


def hairline(slide, x, y, w, color=HAIRLINE, weight=0.9):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Emu(int(weight * 9525)))
    ln.shadow.inherit = False
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln


# ── 슬라이드 셸 ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
BLANK = prs.slide_layouts[6]
_page = {"n": 0}


def new_slide(kind="content", label=None, title=None, subtitle=None, numbered=True):
    s = prs.slides.add_slide(BLANK)
    _page["n"] += 1

    if kind == "content":
        bg = rect(s, 0, 0, SW, SH)
        bg.fill.gradient()
        stops = bg.fill.gradient_stops
        stops[0].color.rgb = GRAD_A
        stops[0].position = 0.0
        stops[1].color.rgb = GRAD_B
        stops[1].position = 1.0
        bg.fill.gradient_angle = 45.0
    elif kind == "blue":
        rect(s, 0, 0, SW, SH, fill=LINK)
    elif kind == "ink":
        rect(s, 0, 0, SW, SH, fill=INK)

    on_dark = kind in ("blue", "ink")
    if label:
        tf = textbox(s, LX, LABEL_Y, W, 0.24)
        para(tf, label, F_MED, 11, ON_PRI_S if on_dark else MUTE, first=True,
             spacing=0.6)
    if title:
        tf = textbox(s, LX, TITLE_Y, W, 0.50)
        para(tf, title, F_BOLD, 30, ON_PRI if on_dark else INK, first=True)
    if subtitle:
        tf = textbox(s, LX, SUB_Y, W, 0.32)
        para(tf, subtitle, F_LIGHT, 14, ON_PRI_S if on_dark else BODY, first=True)

    if numbered:
        tf = textbox(s, LX, FOOTER_Y, 1.0, 0.20)
        para(tf, f"{_page['n']:02d}", F_MED, 10, ON_PRI_S if on_dark else MUTE,
             first=True)
    return s


def bullets(slide, items, x=LX, y=CONTENT_Y, w=W, size=12.5, color=INK,
            family=F_REG, gap=0.30, tag="bullets"):
    """불릿 목록 — 반환값은 마지막 요소 하단 y."""
    cur = y
    for it in items:
        indent = it.get("indent", 0)
        fam = {"b": F_SEMI, "r": F_REG, "l": F_LIGHT}[it.get("w", "r")]
        col = it.get("color", color)
        sz = it.get("size", size)
        marker = "·" if indent else "—"
        tf = textbox(slide, x + 0.14 * indent, cur, w - 0.14 * indent, 0.26)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"{marker}  "
        font(r, fam, sz, LINK if not indent else MUTE)
        r2 = p.add_run(); r2.text = it["t"]
        font(r2, fam, sz, col)
        cur += it.get("gap", gap)
    _check(cur, tag)
    return cur


def table(slide, cols, rows, y=CONTENT_Y, widths=None, size=11.5, hdr_size=11,
          row_h=0.46, tag="table"):
    """하드코딩 없는 하프라인 표 — Rev 7 플랫 스타일."""
    n = len(cols)
    widths = widths or [W / n] * n
    xs, acc = [], LX
    for cw in widths:
        xs.append(acc); acc += cw

    tf = None
    for i, c in enumerate(cols):
        tf = textbox(slide, xs[i], y, widths[i] - 0.10, 0.24)
        para(tf, c, F_SEMI, hdr_size, LINK, first=True)
    hairline(slide, LX, y + 0.28, W, color=ACC_BORDER, weight=1.2)

    cur = y + 0.40
    for r_i, row in enumerate(rows):
        h = row_h
        for i, cell in enumerate(row):
            fam = F_SEMI if (i == 0) else F_REG
            col = INK if (i == 0) else BODY
            tf = textbox(slide, xs[i], cur, widths[i] - 0.10, h,
                         anchor=MSO_ANCHOR.TOP)
            para(tf, cell, fam, size, col, first=True, line=1.15)
        cur += h
        if r_i < len(rows) - 1:
            hairline(slide, LX, cur - 0.06, W)
    _check(cur, tag)
    return cur


OUT = pathlib.Path(__file__).parent / "nyangbti-presentation.pptx"

# ══════════════════════════════════════════════════════════════════════════
# 01 · 표지
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("blue", numbered=False)
tf = textbox(s, LX, 0.42, W, 0.26)
para(tf, "AI 활용 웹콘텐츠 실무역량 과정", F_SEMI, 11, ON_PRI_S, first=True, spacing=1.4)
tf = textbox(s, LX, 1.72, 7.4, 1.0)
para(tf, "냥BTI", F_BLACK, 56, ON_PRI, first=True)
tf = textbox(s, LX, 2.92, 8.2, 0.34)
para(tf, "고양이 성격 유형 검사 웹앱 — 시장조사 · 경쟁사 분석 · 페르소나",
     F_LIGHT, 15, ON_PRI_S, first=True)
hairline(s, LX, 3.52, 2.2, color=ON_PRI_S, weight=1.4)
tf = textbox(s, LX, 3.78, 4.0, 0.80)
para(tf, "2026-08-03", F_LIGHT, 15, ON_PRI_S, first=True, line=1.6)
para(tf, "발표자  srit", F_LIGHT, 15, ON_PRI_S, line=1.6)
for cx, cy, cd in ((8.55, 4.30, 1.55), (9.55, 3.42, 0.95)):
    rect(s, cx - cd / 2, cy - cd / 2, cd, cd, fill=LINK_TINT, shape=MSO_SHAPE.OVAL)

# ══════════════════════════════════════════════════════════════════════════
# 02 · 목차 — flow-steps 4단
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "00 · 목차", "발표 순서", "조사 → 시장 구조 → 빈 자리 → 페르소나")
gap = 0.12
bw = tile(5, gap)
steps = [("01", "조사 방법", "무엇을 어떻게\n검증했는가"),
         ("02", "시장 구조", "세 갈래로\n갈린 시장"),
         ("03", "빈 자리", "아무도\n안 하는 것"),
         ("04", "페르소나", "핵심 사용자\n3명"),
         ("05", "여정 지도", "감정이 꺾이는\n지점")]
BOX_Y, BOX_H = 2.05, 1.62
for i, (num, lab, desc) in enumerate(steps):
    x = LX + i * (bw + gap)
    box = rect(s, x, BOX_Y, bw, BOX_H, fill=LINK)
    tf = box.text_frame
    para(tf, num, F_XBOLD, 16, ON_PRI_S, first=True)
    para(tf, lab, F_SEMI, 12, ON_PRI, space_before=4)
    para(tf, desc, F_LIGHT, 10.5, ON_PRI_S, space_before=4, line=1.25)
    if i < len(steps) - 1:
        ay = BOX_Y + BOX_H / 2
        rect(s, x + bw + 0.028, ay - 0.012, gap - 0.056, 0.024, fill=LINK)
_check(BOX_Y + BOX_H, "02 flow")

# ══════════════════════════════════════════════════════════════════════════
# 03 · 조사 방법 — 2×2 카드
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "01 · 조사 방법", "AI 결과를 그대로 쓰지 않았습니다",
              "조사와 검증을 분리해 각각 독립 수행 후 대조")
cgap = 0.14
cw = tile(2, cgap)
avail = CONTENT_BOT - (CONTENT_Y + 0.08)
ch = (avail - cgap) / 2
cards = [("원칙 1 — 실물 확인", "경쟁사는 스토어·공식 사이트를 직접 열어 확인한 것만 표에 등재"),
         ("원칙 2 — 출처 요구", "모든 행에 검증 근거 표기 (스토어 등재 / 공식 사이트 / 언론 / 논문)"),
         ("원칙 3 — 수치 금지", "1차 출처 없는 다운로드·이용자·시장 규모는 미확인으로 표기"),
         ("원칙 4 — 교차 검증", "조사 담당과 검증 담당을 분리, 대조에서 뒤집힌 항목을 기록")]
for i, (h, b) in enumerate(cards):
    r_i, c_i = divmod(i, 2)
    x = LX + c_i * (cw + cgap)
    y = CONTENT_Y + 0.08 + r_i * (ch + cgap)
    fill = CANVAS if r_i == 0 else CANVAS_S2
    card = rect(s, x, y, cw, ch, fill=fill, line=HAIRLINE)
    rect(s, x, y, 0.05, ch, fill=LINK)
    tf = card.text_frame
    para(tf, h, F_SEMI, 12.5, INK, first=True)
    para(tf, b, F_REG, 11.5, BODY, space_before=6, line=1.25)
_check(CONTENT_Y + 0.08 + 2 * ch + cgap, "03 cards")

# ══════════════════════════════════════════════════════════════════════════
# 04 · 시장 구조 숫자 — metrics-grid
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "02 · 시장 구조", "숫자 두 개가 핵심입니다",
              "확인된 경쟁사 수와, 아무도 하지 않는 것의 개수")
STRIP_Y, STRIP_H = CONTENT_Y + 0.04, 0.92
rect(s, LX, STRIP_Y, W, STRIP_H, fill=ACC_SOFT, line=ACC_BORDER)
metrics = [("8", "직접 경쟁 서비스 (확인됨)"), ("7+", "반려묘 관리 앱 (확인됨)"),
           ("0", "관리 앱의 성향 기능"), ("0", "Play 고양이 성격 앱")]
mcw = W / 4
for i, (fig, lab) in enumerate(metrics):
    x = LX + i * mcw
    tf = textbox(s, x + 0.14, STRIP_Y + 0.14, mcw - 0.28, 0.34)
    para(tf, fig, F_XBOLD, 21, LINK, first=True)
    tf = textbox(s, x + 0.14, STRIP_Y + 0.54, mcw - 0.28, 0.26)
    para(tf, lab, F_MED, 10, MUTE, first=True)
    if i:
        hl = slide_x = x
        ln = rect(s, x, STRIP_Y + 0.16, 0.008, STRIP_H - 0.32, fill=ACC_BORDER)
g2_y = STRIP_Y + STRIP_H + 0.16
g2_h = (CONTENT_BOT - g2_y - cgap) / 2
grid = [("국내 1등 부재", "카테고리에서 먼저 떠오르는 이름이 없음"),
        ("관리 앱은 기록만", "무엇을 했는가는 쌓지만 어떤 아이인가는 다루지 않음"),
        ("국내 테스트는 집사를 측정", "푸망·BLTI·포캣멍센터 — 대상이 우리 고양이가 아님"),
        ("경쟁은 웹에만", "Play 검색 12개 전부 사람용 범용 MBTI 앱")]
for i, (h, b) in enumerate(grid):
    r_i, c_i = divmod(i, 2)
    x = LX + c_i * (cw + cgap)
    y = g2_y + r_i * (g2_h + cgap)
    card = rect(s, x, y, cw, g2_h, fill=CANVAS if r_i == 0 else CANVAS_S2,
                line=HAIRLINE)
    rect(s, x, y, 0.05, g2_h, fill=LINK)
    tf = card.text_frame
    para(tf, h, F_SEMI, 12, INK, first=True)
    para(tf, b, F_REG, 11, BODY, space_before=4, line=1.2)
_check(g2_y + 2 * g2_h + cgap, "04 grid")

# ══════════════════════════════════════════════════════════════════════════
# 05 · 직접 경쟁 — 표 4행
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "02 · 시장 구조", "직접 경쟁 — 4곳으로 압축",
              "전체 8곳 중 전략적으로 의미 있는 4곳")
table(s, ["서비스", "강점", "약점"],
      [["PetMBTI", "국내 유일 한국어 고양이 MBTI + 처방", "기능 산만, 근거 미공개, URL 체계 난잡"],
       ["PurrJung", "정식 한국어 로케일, 논문 인용", "처방 없음, \"검증된 모델\" 과장"],
       ["IDRLabs", "논문 프레임 그대로 구현", "한국어 없음, 권고 제공 거부 명시"],
       ["DBTI", "근거·처방·공유·커뮤니티 완비", "강아지 전용 — 고양이판 미출시"]],
      y=CONTENT_Y + 0.10, widths=[1.70, 3.55, 3.75], row_h=0.62, tag="05 table")

# ══════════════════════════════════════════════════════════════════════════
# 06 · DBTI 벤치마크 — split-panel
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "02 · 시장 구조", "벤치마크 — DBTI",
              "저희가 만들려는 것을 강아지로 이미 만든 곳")
PW = (W - 0.16) / 2
PY, PH = CONTENT_Y + 0.10, CONTENT_BOT - (CONTENT_Y + 0.10)
rect(s, LX, PY, PW, PH, fill=LINK)
tf = textbox(s, LX + 0.22, PY + 0.24, PW - 0.44, 0.72)
para(tf, "카테고리 1등은\n강아지판", F_BLACK, 26, ON_PRI, first=True, line=1.12)
hairline(s, LX + 0.22, PY + 1.22, PW - 0.44, color=ON_PRI_S, weight=1.2)
cur = PY + 1.42
for t in ["22문항 · 5분 · 16유형", "유형별 훈련 · 식이 · 운동 가이드",
          "PDF · 링크 공유, 같은 유형 커뮤니티", "행동 전문가 공동 개발"]:
    tf = textbox(s, LX + 0.22, cur, PW - 0.44, 0.26)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "·  "; font(r, F_REG, 12, ON_PRI_S)
    r2 = p.add_run(); r2.text = t; font(r2, F_REG, 12, ON_PRI)
    cur += 0.30

RX = LX + PW + 0.16
rect(s, RX, PY, PW, PH, fill=CANVAS, line=HAIRLINE)
tf = textbox(s, RX + 0.22, PY + 0.24, PW - 0.44, 0.28)
para(tf, "여기서 배울 것 / 다를 것", F_SEMI, 13, INK, first=True)
hairline(s, RX + 0.22, PY + 0.62, PW - 0.44)
cur = PY + 0.80
for lab, desc in [("배울 것 — 근거·처방·공유 3종 세트", "검사만 주고 끝내지 않는 구조"),
                  ("다를 것 — 고양이 자리는 공석", "FAQ에 고양이는 별도 프레임워크 필요라고 명시"),
                  ("주의 — CBTI 준비 중", "1등이 같은 공백을 노린다고 공표한 상태")]:
    tf = textbox(s, RX + 0.22, cur, PW - 0.44, 0.24)
    para(tf, lab, F_SEMI, 12, LINK, first=True)
    tf = textbox(s, RX + 0.22, cur + 0.26, PW - 0.44, 0.44)
    para(tf, desc, F_REG, 11.5, BODY, first=True, line=1.2)
    cur += 0.86
_check(PY + PH, "06 split")

# ══════════════════════════════════════════════════════════════════════════
# 07 · 학술 근거
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "02 · 시장 구조", "근거는 실재하는 논문에서",
              "Feline Five — PLOS ONE (2017), 반려묘 2,802마리")
end = bullets(s, [
    {"t": "5요인 — 신경성 · 외향성 · 우월성 · 충동성 · 친화성"},
    {"t": "논문이 보호자 실천 권고를 직접 제시 → 처방을 지어내지 않아도 됨"},
    {"t": "국내 경쟁사 중 이를 한국어로 하는 곳 없음", "w": "b"},
], y=CONTENT_Y + 0.08, tag="07 bullets")

PAN_Y = end + 0.14
PAN_H = 1.32
rect(s, LX, PAN_Y, W, PAN_H, fill=ACC_SOFT, line=ACC_BORDER)
tf = textbox(s, LX + 0.20, PAN_Y + 0.16, W - 0.40, 0.26)
para(tf, "단, 그대로 쓰면 안 되는 두 가지", F_SEMI, 12, LINK, first=True)
cur = PAN_Y + 0.48
for t in ["논문이 스스로를 탐색적 연구로 규정 — \"검증된 측정도구\"가 아님",
          "5요인은 연속 점수, 이분 축 아님 → \"Feline Five 기반 16유형\"은 성립하지 않는 문장"]:
    tf = textbox(s, LX + 0.20, cur, W - 0.40, 0.26)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "·  "; font(r, F_REG, 11.5, LINK)
    r2 = p.add_run(); r2.text = t; font(r2, F_REG, 11.5, INK)
    cur += 0.30
tf = textbox(s, LX + 0.20, cur + 0.02, W - 0.40, 0.26)
para(tf, "→ 저희 표기: 문항을 Feline Five에서 참고, 유형화는 재미를 위한 단순화",
     F_SEMI, 11.5, INK, first=True)
_check(PAN_Y + PAN_H, "07 panel")

# ══════════════════════════════════════════════════════════════════════════
# 08 · statement
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("blue")
tf = textbox(s, LX + 0.30, 2.10, W - 0.60, 1.60, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "시장은 관리에 몰려 있고,\n비어 있는 것은 관계 이해입니다.",
     F_BLACK, 34, ON_PRI, first=True, line=1.28)

# ══════════════════════════════════════════════════════════════════════════
# 09 · 차별화 — feature-cards 3-up
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "03 · 빈 자리", "차별화 포인트 3",
              "전부 백엔드 없이 구현 가능한 범위")
fw = tile(3, 0.12)
FY, FH = CONTENT_Y + 0.10, 2.42
feats = [("01", "금지 목록", "유형별 \"이 아이한테 하지 마세요\". 전 경쟁사 미제공 — 처방은 전부 권장 목록뿐"),
         ("02", "근거 + 한계 명시", "논문 참고를 밝히고 탐색적 연구·축 불일치까지 표기. 경쟁사가 과장하는 지점을 정확히 씀"),
         ("03", "백엔드 없이 공유", "유형별 정적 결과 페이지 + og:image → 링크 미리보기가 유형별로 다름")]
for i, (badge, ttl, body) in enumerate(feats):
    x = LX + i * (fw + 0.12)
    card = rect(s, x, FY, fw, FH, fill=CANVAS, line=HAIRLINE)
    bd = rect(s, x + 0.16, FY + 0.18, 0.44, 0.24, fill=ACC_SOFT, line=ACC_BORDER,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tfb = bd.text_frame
    tfb.margin_left = tfb.margin_right = 0
    para(tfb, badge, F_SEMI, 9, LINK, first=True, align=PP_ALIGN.CENTER)
    tf = textbox(s, x + 0.16, FY + 0.54, fw - 0.32, 0.28)
    para(tf, ttl, F_SEMI, 12, INK, first=True)
    hairline(s, x + 0.16, FY + 0.90, fw - 0.32)
    tf = textbox(s, x + 0.16, FY + 1.04, fw - 0.32, FH - 1.20)
    para(tf, body, F_REG, 11.5, BODY, first=True, line=1.28)
tf = textbox(s, LX, FY + FH + 0.14, W, 0.24)
para(tf, "03은 신규 발상 아님 — 1등 플랫폼(푸망) 패턴을 백엔드 없이 정적으로 재현한 것",
     F_MED, 10.5, MUTE, first=True)
_check(FY + FH + 0.38, "09 cards")

# ══════════════════════════════════════════════════════════════════════════
# 10 · divider
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("ink", numbered=True)
tf = textbox(s, LX, 2.28, 6.6, 0.90)
para(tf, "페르소나", F_BLACK, 50, ON_PRI, first=True)
hairline(s, LX, 3.36, 2.2, color=HAIR_STR, weight=1.4)
tf = textbox(s, LX, 3.58, 6.6, 0.30)
para(tf, "핵심 사용자 3명과 주 페르소나 지정", F_LIGHT, 14, HAIR_STR, first=True)
rect(s, 8.55, 3.48, 1.55, 1.55, fill=LINK, shape=MSO_SHAPE.OVAL)

# ══════════════════════════════════════════════════════════════════════════
# 11 · 프로토 페르소나 고지
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "04 · 페르소나", "프로토 페르소나입니다",
              "사용자 인터뷰·설문 없이, 경쟁사 분석과 추론으로 만든 가설")
end = table(s, ["항목", "내용"],
            [["유형", "프로토 페르소나 — NN/g가 인정하는 3유형 중 하나"],
             ["근거", "경쟁사 18행 + 논문 3편 + 검증 실패 13건"],
             ["용도", "설계 초점 도구 — 데이터로 주장하지 않음"],
             ["검증 계획", "실제 집사 5명 롤링 인터뷰"]],
            y=CONTENT_Y + 0.06, widths=[1.70, 7.30], row_h=0.42, tag="11 table")
PAN_H = 0.86
rect(s, LX, end + 0.12, W, PAN_H, fill=ACC_SOFT, line=ACC_BORDER)
tf = textbox(s, LX + 0.20, end + 0.26, W - 0.40, 0.26)
para(tf, "가장 약한 가정 — \"집사가 자기 고양이 행동을 문항으로 답할 만큼 관찰하고 있다\"",
     F_SEMI, 11.5, INK, first=True)
tf = textbox(s, LX + 0.20, end + 0.56, W - 0.40, 0.26)
para(tf, "틀리면 검사 구조 전체가 무너짐 → 대응으로 \"모르겠음\" 선택지와 관찰형 문항 도입",
     F_REG, 11, BODY, first=True)
_check(end + 0.12 + PAN_H, "11 panel")

# ══════════════════════════════════════════════════════════════════════════
# 12 · 행동 축
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "04 · 페르소나", "인구통계가 아니라 행동으로 나눴습니다",
              "축을 밝히지 않으면 고정관념으로 읽힙니다")
end = table(s, ["행동 축", "A 초보", "B 맞벌이", "C 다묘"],
            [["해석 공백", "판단 기준 자체가 없음", "기준은 있으나 시간 없음", "비교 축이 없음"],
             ["떨어져 있는 시간", "짧음 — 과잉 접촉", "길다 — 평일 10시간 이상", "중간 — 3마리로 분산"],
             ["구매 계기", "추천 전부 시도 → 실패 누적", "1종만 정확히", "신규 구매 아니라 배분"],
             ["재검사", "2~4주 내 의향 높음", "1회성", "합사·서열 변동 시"]],
            y=CONTENT_Y + 0.06, widths=[1.86, 2.44, 2.44, 2.26], size=11,
            row_h=0.58, tag="12 table")
tf = textbox(s, LX, end + 0.12, W, 0.26)
para(tf, "자체 검사 — 세 명의 인구통계를 맞바꿔도 필요 기능이 바뀌지 않으면 실패. 통과 확인",
     F_SEMI, 11, LINK, first=True)
_check(end + 0.38, "12 note")

# ══════════════════════════════════════════════════════════════════════════
# 여정맵 읽는 법 — 페르소나 짝 앞에 한 장
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "05 · 여정 지도", "고객 여정 지도 — 무엇이고 무엇이 아닌가",
              "7단계는 상세판 · 리서치 없는 가설 저니맵(assumption map)")
end = table(s, ["필수 요소 (NN/g)", "우리 문서에서의 위치", "상태"],
            [["Actor — 맵당 1명", "페르소나 3인 각각 별도 맵", "충족"],
             ["Scenario + 기대", "페르소나별 기대 한 줄", "충족"],
             ["Journey Phases", "7단계 — 강사 배포 양식 그대로", "충족"],
             ["Actions · Mindsets · Emotions", "생각·행동 / 감정 / Pain Point", "충족"],
             ["Opportunities + 담당 + 지표", "기회 요소 (담당·확인 방법은 뒤에)", "보완"]],
            y=CONTENT_Y + 0.02, widths=[2.55, 4.05, 2.40], size=11, hdr_size=10.5,
            row_h=0.38, tag="jm-intro table")
cur = end + 0.10
tf = textbox(s, LX, cur, W, 0.24)
para(tf, "숨기지 않고 먼저 밝히는 3가지", F_SEMI, 11.5, LINK, first=True)
for t in ["인터뷰·설문·애널리틱스 0건 — 근거는 경쟁사 18행 + 논문 3편 + 행동 축",
          "감정선은 관찰 결과가 아니라 설계 가설 — 최저점·회복 지점을 우리가 지정",
          "감정 척도  −2 매우 부정 · −1 부정 · 0 중립 · +1 긍정 · +2 매우 긍정"]:
    cur += 0.23
    tf = textbox(s, LX, cur, W, 0.22)
    pp = tf.paragraphs[0]
    r = pp.add_run(); r.text = "·  "; font(r, F_REG, 11, MUTE)
    r2 = pp.add_run(); r2.text = t; font(r2, F_REG, 11, BODY)
_check(cur + 0.22, "jm-intro labels")

# ══════════════════════════════════════════════════════════════════════════
# 페르소나 + 그 사람의 여정맵을 짝으로 — 1-6 강사 지시
#   "위에 김소연 퍼소나, 그 밑에 김소연 여정맵" — 한 사람 것을 붙여 둬야 파악됨
#   순서는 등급순 A(primary) → C(secondary) → B(supplemental)
# ══════════════════════════════════════════════════════════════════════════
personas = [
    dict(tier="주 페르소나 (primary)", name="오세진", cap="초보 집사",
         prof=["20대 후반 · 1인 가구", "수도권 도심 원룸 월세",
               "코숏 1살 미만 · 함께 4개월", "퇴근 후 계속 붙어 있음 — 과잉 접촉",
               "월 지출 10~20만원대"],
         motive="우리 고양이 행동을 하나도 해석 못 하겠음",
         need="이 행동이 무슨 뜻인지 알고 싶음",
         tobe="유형을 알고 접촉 강도를 조절, 하지 말아야 할 행동을 구분",
         feat="유형 서사 · 금지 목록 · \"모르겠음\" 선택지",
         img="1st-crop.png",
         expect="\"왜 그러는지 쉽게 알고 싶어\"",
         aha="금지 목록",
         journey=[
             ["인지", "커뮤니티에서 결과 카드를 봄", "0 기대", "검색 결과가 전부 일반론", "결과 카드가 유입 경로"],
             ["유입", "\"가입 요구하면 나간다\"", "−1 경계", "가입·권한 요구에 민감", "회원가입 없음 첫 화면 명시"],
             ["온보딩", "\"내가 답할 수 있나\"", "−1 불안", "관찰 안 한 항목이 두려움", "최근 일주일 + 모르겠음 예고"],
             ["테스트 진행", "애매하면 모르겠음 선택", "0 몰입", "비교 기준 자체가 없음", "막힘 없이 끝까지"],
             ["결과 확인", "\"그래서 뭘 해야 하지?\"", "−2 아쉬움", "설명만으론 행동이 안 바뀜", "상단에 처방 진입점"],
             ["케어 가이드", "금지 목록에서 원인 발견", "+2 납득", "권장만으론 구분이 안 됨", "금지 목록 — 결정적 순간"],
             ["재방문", "2~4주 후 재검사·비교", "+1 기대", "이전 결과 없으면 변화 모름", "검사일 기록 · 재검사 안내"],
         ]),
    dict(tier="2차 페르소나 (secondary)", name="임현석", cap="다묘가정 집사",
         prof=["40대 초반 · 배우자 + 자녀 1", "지방 중소도시 구축 아파트 자가",
               "코숏 8살 / 코숏 4살 / 페르시안 1살", "합사 이력 3회 상이 — 갈등 원인 불명",
               "월 지출 40~60만원대"],
         motive="갈등의 출처가 누구인지 판정이 안 됨",
         need="세 마리 성향을 나란히 비교해 보고 싶음",
         tobe="유형 조합을 근거로 자원 분리와 배분을 조정",
         feat="다묘 비교(localStorage) · 유형 조합별 자원 분리 권고",
         img="3rd-crop.png",
         expect="\"셋 중 누가 문제인지 알고 싶어\"",
         aha="유형 조합별 자원 분리",
         journey=[
             ["인지", "다묘 합사 갈등 글에서 발견", "−1 반신반의", "다묘 자료가 거의 없음", "다묘 키워드 진입 경로"],
             ["유입", "\"세 마리 다 해야 하나\"", "−1 망설임", "반복 입력이 진입 장벽", "한 마리만 먼저 해도 됨"],
             ["온보딩", "개체별로 따로 답해야 함", "−1 부담", "공통 정보를 3회 재입력", "가구 정보는 재사용"],
             ["테스트 진행", "3회 반복 응답", "−1 피로", "3회째 집중력 저하", "마리별 저장 · 재개"],
             ["결과 확인", "개체별 유형만 나옴", "−2 답답함", "나란히 볼 수가 없음", "비교 화면 진입점"],
             ["케어 가이드", "유형 조합 자원 분리 확인", "+2 납득", "일반 조언은 \"천천히 합사\"", "조합별 권고 — 결정적 순간"],
             ["재방문", "합사·서열 변동 때마다", "0 중립", "변동 시점을 직접 기억", "개체 추가 시 비교 갱신"],
         ]),
    dict(tier="보조 페르소나 (supplemental)", name="문서연", cap="맞벌이 집사",
         prof=["30대 중반 · 부부 2인", "경기권 신도시 아파트 전세",
               "러시안 블루 3살 · 함께 약 3년", "평일 10시간 이상 부재, 밤 20~30분 놀이",
               "월 지출 20~30만원대"],
         motive="짧은 시간이라도 제대로 해주고 싶음",
         need="실패 없이 딱 하나만 정확히 사고 싶음",
         tobe="유형별 장난감 우선순위 1종을 확인하고 구매",
         feat="처방 3블록 · 우선순위 1종 표기 · 배우자 공유 링크",
         img="2nd-crop.png",
         expect="\"짧은 시간에 제대로 해주고 싶어\"",
         aha="우선순위 1종",
         journey=[
             ["인지", "배우자 링크 · SNS 노출", "0 관심", "탐색할 시간이 없음", "미리보기에 소요 시간"],
             ["유입", "밤에 짧게 확인", "−1 조급", "로딩·가입이 이탈 지점", "즉시 시작 · 가입 없음"],
             ["온보딩", "\"3분이면 지금 되네\"", "+1 안심", "시간 불명확하면 미룸", "소요 시간 버튼 옆 명시"],
             ["테스트 진행", "한 번에 몰아서 응답", "0 무난", "끊기면 다시 안 옴", "짧은 길이 유지"],
             ["결과 확인", "\"그래서 뭘 사면 되지?\"", "−2 미흡", "유형 서사는 효용이 낮음", "구매 판단으로 직행"],
             ["케어 가이드", "우선순위 1종 확인", "+2 결정", "3~5개면 고민이 늘어남", "\"이것부터\" — 결정적 순간"],
             ["재방문", "1회성으로 끝날 가능성", "0 중립", "재방문 설계로 안 잡힘", "배우자 공유로 확산 대체"],
         ]),
]
for p in personas:
    s = new_slide("content", "04 · 페르소나", f"{p['name']} — {p['cap']}", p["tier"])
    PY2, PH2 = CONTENT_Y + 0.10, CONTENT_BOT - (CONTENT_Y + 0.10)
    rect(s, LX, PY2, PW, PH2, fill=LINK)
    # 이미지 박스 3.98 x 1.45 (약 2.75:1 와이드) — 1.06"는 3.75:1 띠라서
    # 인물+고양이+실내가 함께 담기지 않는다. 생성은 16:9로 하고 여기서 크롭.
    # 크롭본은 1408x513 = 2.745:1 로 박스 비율과 정확히 일치 — 왜곡 없음
    imgp = pathlib.Path(__file__).parent.parent / "assets" / p["img"]
    s.shapes.add_picture(str(imgp), Inches(LX + 0.22), Inches(PY2 + 0.20),
                         width=Inches(PW - 0.44), height=Inches(1.45))
    cur = PY2 + 1.81
    for t in p["prof"]:
        tf = textbox(s, LX + 0.22, cur, PW - 0.44, 0.26)
        pp = tf.paragraphs[0]
        r = pp.add_run(); r.text = "·  "; font(r, F_REG, 11.5, ON_PRI_S)
        r2 = pp.add_run(); r2.text = t; font(r2, F_REG, 11.5, ON_PRI)
        cur += 0.29
    rect(s, RX, PY2, PW, PH2, fill=CANVAS, line=HAIRLINE)
    cur = PY2 + 0.22
    # 0.78 스텝 — 0.88이면 필요 기능 블록이 5.38"까지 내려가 푸터 구역을 침범한다
    for lab, val in [("동기 (AS-IS)", p["motive"]), ("니즈 (AS-IS)", p["need"]),
                     ("TO-BE", p["tobe"])]:
        tf = textbox(s, RX + 0.22, cur, PW - 0.44, 0.24)
        para(tf, lab, F_SEMI, 11.5, LINK, first=True)
        tf = textbox(s, RX + 0.22, cur + 0.25, PW - 0.44, 0.48)
        para(tf, val, F_REG, 11.5, INK, first=True, line=1.20)
        cur += 0.78
    hairline(s, RX + 0.22, cur - 0.06, PW - 0.44)
    tf = textbox(s, RX + 0.22, cur + 0.02, PW - 0.44, 0.20)
    para(tf, "필요 기능", F_SEMI, 11, MUTE, first=True)
    tf = textbox(s, RX + 0.22, cur + 0.24, PW - 0.44, 0.46)
    para(tf, p["feat"], F_REG, 11, BODY, first=True, line=1.18)
    _check(cur + 0.70, f"persona {p['name']} 필요기능")

    # --- 바로 다음 장에 이 사람의 여정맵 (강사 지시: 짝으로 붙여 둘 것) ---
    s = new_slide("content", "05 · 여정 지도", f"{p['name']}의 여정맵 — 7단계",
                  f"기대 — {p['expect']}")
    end = table(s, ["단계", "생각·행동", "감정", "Pain Point", "기회 요소"],
                p["journey"], y=CONTENT_Y + 0.02,
                widths=[1.00, 2.25, 0.90, 2.35, 2.50], size=9.5, hdr_size=9.5,
                row_h=0.38, tag=f"journey {p['name']}")
    tf = textbox(s, LX, end + 0.08, W, 0.22)
    pp = tf.paragraphs[0]
    r = pp.add_run(); r.text = "최저점 결과 확인 (−2)  →  회복 지점 케어 가이드 (+2)   ·   회복시킨 것 = "
    font(r, F_REG, 10.5, BODY)
    r2 = pp.add_run(); r2.text = p["aha"]; font(r2, F_SEMI, 10.5, LINK)
    _check(end + 0.30, f"journey note {p['name']}")

# --- 저니맵 ④ 감정 곡선 (실제 꺾은선) ------------------------------------
s = new_slide("content", "05 · 여정 지도", "감정 곡선 한눈에 보기",
              "앞 세 장의 감정 열을 하나의 선으로 — 어디서 꺾이는지 바로 보이도록")
STAGES = ["인지", "유입", "온보딩", "테스트", "결과", "케어", "재방문"]
SERIES = [("A 초보", [0, -1, -1, 0, -2, 2, 1], LINK, 2.25),
          ("C 다묘", [-1, -1, -1, -1, -2, 2, 0], INK, 1.5),
          ("B 맞벌이", [0, -1, 1, 0, -2, 2, 0], MUTE, 1.5)]
# 감정은 라벨이 아니라 선으로 그려야 한다는 NN/g 요건 — 실제 꺾은선으로 렌더
PX, PY3, PWID, PHGT = LX + 0.62, CONTENT_Y + 0.10, W - 0.62, 1.86
for v in (2, 1, 0, -1, -2):
    gy = PY3 + (2 - v) / 4 * PHGT
    hairline(s, PX, gy, PWID, color=ACC_BORDER if v == 0 else HAIRLINE,
             weight=1.4 if v == 0 else 0.9)
    tf = textbox(s, LX, gy - 0.10, 0.55, 0.20, align=PP_ALIGN.RIGHT)
    para(tf, f"{v:+d}" if v else "0", F_MED, 9.5, MUTE, first=True)
step_x = PWID / (len(STAGES) - 1)
for i, nm in enumerate(STAGES):
    tf = textbox(s, PX + i * step_x - 0.45, PY3 + PHGT + 0.06, 0.90, 0.20,
                 align=PP_ALIGN.CENTER)
    para(tf, nm, F_MED, 9.5, BODY, first=True)
for label, vals, col, thick in SERIES:
    pts = [(PX + i * step_x, PY3 + (2 - v) / 4 * PHGT) for i, v in enumerate(vals)]
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
        cn.line.color.rgb = col
        cn.line.width = Pt(thick)
    for (x, y) in pts:
        d = 0.075 if col == LINK else 0.058
        rect(s, x - d / 2, y - d / 2, d, d, fill=col, shape=MSO_SHAPE.OVAL)
lx = PX
for label, vals, col, thick in SERIES:
    rect(s, lx, PY3 - 0.16, 0.22, 0.045, fill=col)
    tf = textbox(s, lx + 0.28, PY3 - 0.26, 1.15, 0.20)
    para(tf, label, F_SEMI, 10, col if col != MUTE else BODY, first=True)
    lx += 1.55
cur = PY3 + PHGT + 0.34
for t, w in [("세 사람 모두 최저점은 결과 확인, 최고점은 케어 가이드 확인 — 꺾이는 단계는 동일", "b"),
             ("회복시킨 것이 다름 — A는 금지 목록 · C는 유형 조합별 자원 분리 · B는 우선순위 1종", "b"),
             ("결과 확인까지는 경쟁사에서 겪는 현재 상태 / 케어 가이드부터는 아직 없는 화면 = 가설", "r"),
             ("C는 진입 구간이 내내 −1 (세 마리 반복 입력 부담) · B는 재방문이 0에서 멈춤 (1회성)", "r")]:
    tf = textbox(s, LX, cur, W, 0.22)
    pp = tf.paragraphs[0]
    r = pp.add_run(); r.text = "·  "; font(r, F_REG, 10.5, LINK if w == "b" else MUTE)
    r2 = pp.add_run(); r2.text = t
    font(r2, F_SEMI if w == "b" else F_REG, 10.5, INK if w == "b" else BODY)
    cur += 0.23
_check(cur, "jm4 bullets")

# ══════════════════════════════════════════════════════════════════════════
# 16 · 주 페르소나 지정 — flow-steps 3단
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "04 · 페르소나", "주 페르소나 지정 — 소거법",
              "인터페이스가 하나이므로 주 페르소나는 1명")
sw3 = tile(3, 0.12)
SY, SHb = CONTENT_Y + 0.10, 1.52
elim = [("01", "B를 타깃하면", "문항 최소 + 장난감 한 장으로 수렴\n→ A는 입력조차 못 함"),
        ("02", "C를 타깃하면", "비교 화면이 한 마리를 읽을 줄 안다를 전제\n→ A는 진입 과부하"),
        ("03", "A를 타깃하면", "B는 처방·짧은 소요로 최소 만족\nC는 비교 화면 하나 추가로 만족")]
for i, (num, lab, desc) in enumerate(elim):
    x = LX + i * (sw3 + 0.12)
    box = rect(s, x, SY, sw3, SHb, fill=LINK)
    tf = box.text_frame
    para(tf, num, F_XBOLD, 16, ON_PRI_S, first=True)
    para(tf, lab, F_SEMI, 12, ON_PRI, space_before=4)
    para(tf, desc, F_LIGHT, 10.5, ON_PRI_S, space_before=4, line=1.25)
    if i < 2:
        ay = SY + SHb / 2
        rect(s, x + sw3 + 0.028, ay - 0.012, 0.12 - 0.056, 0.024, fill=LINK)
cur = SY + SHb + 0.16
tf = textbox(s, LX, cur, W, 0.26)
para(tf, "주 페르소나 = A 오세진 · 가장 만족시키기 어려운 사람이기도 함",
     F_SEMI, 12.5, LINK, first=True)
for t in ["C = secondary — 다묘 비교 1개만, 첫 화면 노출 금지",
          "B = supplemental — 주 페르소나 설계가 이미 커버. B 전용 기능 추가하지 않음"]:
    cur += 0.28
    tf = textbox(s, LX, cur, W, 0.26)
    pp = tf.paragraphs[0]
    r = pp.add_run(); r.text = "·  "; font(r, F_REG, 11.5, MUTE)
    r2 = pp.add_run(); r2.text = t; font(r2, F_REG, 11.5, BODY)
_check(cur + 0.26, "16 concl")

# ══════════════════════════════════════════════════════════════════════════
# 17 · 공통 니즈 → 기능 — keyword-chips
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "04 · 페르소나", "공통 니즈에서 기능이 나옵니다",
              "셋 다 원하는 것 — 이유는 각각 다름")
SUM_Y, SUM_H = CONTENT_Y + 0.06, 0.90
rect(s, LX, SUM_Y, W, SUM_H, fill=ACC_SOFT, line=ACC_BORDER)
tf = textbox(s, LX + 0.20, SUM_Y + 0.14, W - 0.40, 0.62)
para(tf, "일반론이 아닌 우리 고양이 한 마리의 개별 판독 · 결과가 곧 다음 행동 · "
         "해도 되는 것과 안 되는 것의 구분 · 3분 안팎 · 회원가입 없이 즉시",
     F_REG, 11.5, INK, first=True, line=1.32)
chw = tile(4, 0.12)
CH_Y = SUM_Y + SUM_H + 0.18
for i, c in enumerate(["16문항 이내", "금지 목록", "처방 3블록", "결과 카드 공유"]):
    x = LX + i * (chw + 0.12)
    chip = rect(s, x, CH_Y, chw, 0.50, fill=LINK,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    tfc = chip.text_frame
    tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tfc, c, F_SEMI, 12, ON_PRI, first=True, align=PP_ALIGN.CENTER)
cur = CH_Y + 0.50 + 0.20
colw = (W - 0.20) / 2
for i, (h, items) in enumerate([
        ("채택", "관찰형 문항 · \"모르겠음\" 선택지 · 근거 표기 블록 · 유형별 정적 결과 URL · 다묘 비교"),
        ("범위 외 — 서버·인증 필요", "참여 수 통계 · 커뮤니티 · 기기 간 동기화 · 사진 AI 분석")]):
    x = LX + i * (colw + 0.20)
    tf = textbox(s, x, cur, colw, 0.24)
    para(tf, h, F_SEMI, 11.5, LINK if i == 0 else MUTE, first=True)
    tf = textbox(s, x, cur + 0.26, colw, 0.62)
    para(tf, items, F_REG, 11, BODY, first=True, line=1.25)
_check(cur + 0.88, "17 cols")

# ══════════════════════════════════════════════════════════════════════════
# 여정 지도 마무리 2장 — 감정 곡선 / 기능 요구
# (개별 여정맵은 각 페르소나 바로 뒤에 배치 — 1-6 강사 지시)
# 내용 출처: scratchpad/journey-map-slides.md (1차 출처 조사 기반)
# ══════════════════════════════════════════════════════════════════════════

# --- 저니맵 ⑤ 기능 요구 + 담당·확인 방법 ----------------------------------
s = new_slide("content", "05 · 여정 지도", "저니맵 → 기능 요구 — 담당과 확인 방법까지",
              "기회 요소를 화면으로 바꾼 목록. IA·유저 플로우가 이 표를 그대로 소비한다")
end = table(s, ["기능", "근거가 된 칸", "담당", "확인 방법"],
            [["첫 화면에 회원가입 없음 · 소요 시간 표기", "A 유입 · B 온보딩", "공동",
              "첫 화면에 두 문구가 보이는가"],
             ["검사 전 안내 — 최근 일주일 기준 · 모르겠음 있음", "A 온보딩", "문항",
              "안내 화면이 존재하는가"],
             ["결과 화면 상단에 처방 진입점", "A·B·C 결과 확인", "결과",
              "유형 서사보다 위에 있는가"],
             ["유형별 금지 목록 블록", "A 케어 가이드", "콘텐츠",
              "16유형 전부 2~3개씩 작성됐는가"],
             ["추천 장난감 우선순위 1종 표기", "B 케어 가이드", "콘텐츠",
              "1종이 시각적으로 먼저 보이는가"],
             ["결과 화면에서 비교 화면 진입점", "C 결과 확인", "결과",
              "2마리 이상일 때만 활성화되는가"]],
            y=CONTENT_Y + 0.04, widths=[3.30, 2.10, 0.80, 2.80], size=10,
            hdr_size=10, row_h=0.40, tag="jm5 table")
tf = textbox(s, LX, end + 0.08, W, 0.22)
para(tf, "「확인 방법」 열이 곧 유저 플로우의 분기 조건 — 저니맵은 감정, 플로우는 화면·시스템 응답",
     F_SEMI, 10.5, LINK, first=True)
tf = textbox(s, LX, end + 0.30, W, 0.22)
para(tf, "범위 외 — 재방문 알림(푸시·메일)은 서버 필요 → localStorage 기반 지난 검사일 안내로 대체",
     F_REG, 10, BODY, first=True)
_check(end + 0.52, "jm5 notes")

# ══════════════════════════════════════════════════════════════════════════
# 18 · AI 한계
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "05 · AI 한계", "AI가 틀린 지점을 기록했습니다",
              "조사 과정에서 실제로 발생한 오류")
end = table(s, ["유형", "사례", "확인 방법"],
            [["없는 서비스", "후보 2건이 404 (스토어 등재·테스트 링크)", "URL 직접 접속"],
             ["종료된 서비스", "펫닥 — 양대 스토어에서 삭제", "스토어 API 응답 확인"],
             ["수치 과장", "\"학습 데이터 25만\" → 실제 약 1만", "언론 기사·공식 사이트 대조"],
             ["논문 ID 오류", "검색 요약이 잘못된 PubMed ID 제시", "논문 원문 개별 열람"]],
            y=CONTENT_Y + 0.06, widths=[1.70, 4.60, 2.70], size=11,
            row_h=0.50, tag="18 table")
tf = textbox(s, LX, end + 0.14, W, 0.26)
para(tf, "배운 것 — 출처를 요구하면 거짓이 줄지만, 요구만으로는 부족. 1차 출처를 하나씩 직접 열람",
     F_SEMI, 11.5, LINK, first=True)
_check(end + 0.40, "18 note")

# ══════════════════════════════════════════════════════════════════════════
# 19 · 다음 단계
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "06 · 다음 단계", "다음은 화면 구조와 플로우",
              "여정맵에서 찾은 막히는 지점을 화면으로 바꾸는 순서")
end = bullets(s, [
    {"t": "완료 — 주제 확인 · 시장조사 · 경쟁사 분석 · 페르소나 · 여정 지도 · 유저 시나리오", "w": "b"},
    {"t": "다음 — IA(화면 구조도) → 유저 플로우 → 와이어프레임"},
    {"t": "이후 — 무드보드 → 화면 디자인 → 코딩"},
], y=CONTENT_Y + 0.10, tag="19 bullets")
PAN_H = 1.10
rect(s, LX, end + 0.16, W, PAN_H, fill=CANVAS, line=HAIRLINE)
rect(s, LX, end + 0.16, 0.05, PAN_H, fill=LINK)
tf = textbox(s, LX + 0.22, end + 0.30, W - 0.44, 0.26)
para(tf, "남은 결정 2건", F_SEMI, 12, INK, first=True)
cur = end + 0.58
for t in ["유형 체계 — 4축 16유형 유지(강사 예제 형식) + 축 불일치 명시",
          "최종 명칭 — \"냥BTI\"는 2022년 굿즈 캠페인에 사용된 이력 있음"]:
    tf = textbox(s, LX + 0.22, cur, W - 0.44, 0.26)
    pp = tf.paragraphs[0]
    r = pp.add_run(); r.text = "·  "; font(r, F_REG, 11, LINK)
    r2 = pp.add_run(); r2.text = t; font(r2, F_REG, 11, BODY)
    cur += 0.26
_check(end + 0.16 + PAN_H, "19 panel")

# ══════════════════════════════════════════════════════════════════════════
# 20 · 마지막 장 — 페이지 번호 없음
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("blue", numbered=False)
tf = textbox(s, LX, 0.42, W, 0.26)
para(tf, "AI 활용 웹콘텐츠 실무역량 과정", F_SEMI, 11, ON_PRI_S, first=True, spacing=1.4)
cur = 1.18
for t in ["시장은 관리 중심 — 관계 이해가 비어 있음",
          "주 페르소나 1명으로 좁히고, 나머지는 등급으로 정리",
          "근거는 밝히되 한계까지 같이 적음"]:
    tf = textbox(s, LX, cur, W, 0.28)
    pp = tf.paragraphs[0]
    r = pp.add_run(); r.text = "·  "; font(r, F_REG, 12.5, ON_PRI_S)
    r2 = pp.add_run(); r2.text = t; font(r2, F_REG, 12.5, ON_PRI)
    cur += 0.34
tf = textbox(s, LX, 2.70, 7.0, 0.90)
para(tf, "감사합니다", F_BLACK, 50, ON_PRI, first=True)
rect(s, 0, 4.92, SW, 0.705, fill=LINK_TINT)
tf = textbox(s, LX, 5.14, W, 0.26)
para(tf, "AI 활용 웹콘텐츠 실무역량 과정 · 2026-08-03 · srit",
     F_MED, 10.5, ON_PRI_S, first=True)

# ── 저장 + 검증 ────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"saved: {OUT}  ({OUT.stat().st_size:,} bytes)")
print(f"slides: {len(prs.slides._sldIdLst)}")
print(f"tile widths — 4@0.12={tile(4,0.12):.4f}\"  3@0.12={tile(3,0.12):.4f}\"  "
      f"2@0.14={tile(2,0.14):.4f}\"")

# 저장 후 실제 도형을 전부 훑어 재검증.
# 빌드 중 _check()는 패널 바깥 테두리만 봤기 때문에 패널 *안쪽* 텍스트박스가
# 5.08"를 넘는 것을 놓쳤다 (페르소나 3장의 "필요 기능" 블록이 5.38"까지 내려갔음).
# 이 사후 검증이 그 부류의 실수를 다시 놓치지 않게 하는 안전장치다.
EMU_IN = 914400
CLOSING_IDX = len(prs.slides._sldIdLst)
hard = []
for idx, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        if sh.top is None or sh.width is None or sh.height is None:
            continue
        top, bot, wid = sh.top / EMU_IN, (sh.top + sh.height) / EMU_IN, sh.width / EMU_IN
        if wid > 9.9:                      # 전면 배경 / 마감 스트립
            continue
        if abs(top - FOOTER_Y) < 0.02:     # 푸터 행 — 규격상 정상
            continue
        if idx == CLOSING_IDX and top > 5.0:   # Rev 7: 제출 정보 스트립은 하단 밀착
            continue
        if bot > CONTENT_BOT + 1e-3:
            txt = sh.text_frame.text[:30].replace("\n", "/") if sh.has_text_frame else ""
            hard.append(f"slide {idx}: bottom {bot:.3f}\" — {txt}")

if _violations or hard:
    print(f"\n!!! CONTENT_BOT violations — build-time {len(_violations)}, post-save {len(hard)}")
    for v in _violations + hard:
        print("   ", v)
    raise SystemExit(1)
print(f"CONTENT_BOT {CONTENT_BOT}\" — verified across every shape on all "
      f"{CLOSING_IDX} slides")
