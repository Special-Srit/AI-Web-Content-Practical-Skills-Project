#!/usr/bin/env python3
"""Music Diary 발표 자료 빌더 — Warm Vinyl Deck v1 (색·레이아웃 [팀 확정 필요]).

내용 출처: team-project/docs/presentation-outline.md (2026-08-06 sol 검토 반영본)
구조 참고: nyangbti/presentation/build_deck.py (Rev 7) — 앵커·헬퍼 함수는 그대로
          가져오되 색상 토큰만 Warm Vinyl로 교체. Rev 7 자체를 상속하지 않음.
실행:  <venv>/bin/python build_deck.py

이번 패스에서 실제로 다른 점 (14-deck-tooling-research.md 권고 반영):
  - measure_kr() — Pillow 기반 사전 측정. table()이 고정 row_h 대신 실제
    한국어 텍스트의 예상 줄바꿈 높이를 먼저 계산해 오버플로를 미리 잡는다
    (§4 "adopt" 권고). Rev 7의 CONTENT_BOT 단일 assert보다 한 단계 앞선 검사.
  - 팀 로고·아이콘 세트가 없어 아이콘 배지는 이번 패스에서 생략 (§2 "필요할 때만"
    권고 — 배지 5개 이상 필요해지면 그때 cairosvg 경로 추가)
  - 부록 슬라이드는 이번 패스에서 만들지 않는다 — 본문 8슬라이드만 (outline.md
    "부록" 절 내용은 발표 시 구두 백업으로, 카드 프린트나 Q&A 화면 전환으로 대응)

정정 (2026-08-06, sol 시니어 디자인 검토 + 16-great-ppt-design-research.md 반영) —
첫 렌더(preview/*.png)에서 실제로 잡힌 문제와 그 수정:
  1. 본문·표·캡션이 인터페이스 텍스트 크기(10~12.5pt)였다 — 발표 투사 기준
     18pt 미만은 피하라는 Microsoft 권고 위반. 전부 18pt 기준으로 올림
     (TYPE 토큰 딕셔너리로 통일, §타이포그래피).
  2. ACCENT·MUTE·HAIRLINE이 CANVAS 위에서 WCAG 4.5:1 대비 기준에 못 미쳤다
     (계산: ACCENT 4.385:1, MUTE 3.84:1, HAIRLINE 1.32:1). 셋 다 더 어둡게
     조정 — sol이 실제 계산한 교체값 사용.
  3. 06번 슬라이드의 다이아몬드가 제목·부제 텍스트 박스와 겹쳐 대비가
     무너졌다(흰 텍스트 위에 옅은 다이아몬드가 곧바로 걸침). 다이아몬드를
     우상단 모서리로 옮겨 텍스트 안전 영역(TEXT_SAFE) 밖으로 뺌.
  4. 짧은 표(02·07번)가 콘텐츠 영역 위쪽 1/3만 채우고 하단이 비어 보였다 —
     행 높이를 억지로 늘리는 대신 표 그룹 자체를 세로 중앙에 배치
     (16-great-ppt-design-research.md §1 권고).
  5. 표 열 너비가 실제 텍스트 길이에 안 맞았다(02·03·05·07) — sol이 각 셀의
     실제 글자 수를 세어 다시 계산한 비율로 교체.
  6. 04번 흐름 카드가 내부 텍스트를 위쪽에만 몰아 아래가 허전했다 — 카드
     텍스트를 세로 중앙 정렬, 연결 막대도 눈에 띄게 키움.
  7. 08번의 팀 사진 자리가 왼쪽에 몰려 오른쪽 절반이 비어 있었다 — 9인치
     콘텐츠 폭 전체에 3등분으로 재배치. 대괄호 placeholder 문구 자체는
     유지(팀 역할·사진은 아직 실제로 미정 — 이건 레이아웃 버그이지 콘텐츠
     결정이 아니므로 여기서 지어내지 않음).
  8. 03·05번의 정정 각주를 캔버스 위 10pt 캡션 대신 발표자 노트로 옮김 —
     sol 권고(글자를 줄이지 말고 speaker notes로 옮기라는 지적).

Warm Vinyl 하드 규칙 (팀 승인 전 — build_deck.py 재실행으로 언제든 바꿀 수 있음):
  - 콘텐츠 상한 CONTENT_BOT = 5.08" (냥BTI와 동일 앵커, 16:9 마스터 재사용)
  - 등폭 타일 itemW = (W - (N-1)*gap) / N
  - 페이지 번호 2자리 영패딩, 표지/마지막 장 제외
  - 폰트는 Paperlogy Filled 계열 (냥BTI와 동일 — 15-deck-moodboard.md가 미정으로
    남긴 항목을 이번 스크립트에서 "재사용"으로 확정)
  - 이모지 금지
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont
import pathlib

# ── 색상 토큰 (Warm Vinyl Deck v1 — 값(안), 15-deck-moodboard.md §색상 방향) ──
# ACCENT·MUTE·HAIRLINE은 2026-08-06 sol 검토가 실측한 WCAG 대비로 교체됨 —
# 원래 값(#C1502E·#8A7A6D·#E3D7C9)은 CANVAS 위에서 각각 4.385:1·3.84:1·1.32:1로
# 4.5:1 기준에 못 미쳤다. 아래는 sol이 다시 계산해 확인한 값.
CANVAS      = RGBColor(0xFB, 0xF6, 0xEF)   # warm ivory
INK         = RGBColor(0x24, 0x1C, 0x16)   # warm near-black
ACCENT      = RGBColor(0xAC, 0x41, 0x25)   # terracotta/rust, 진하게 — CANVAS 위 5.52:1
ACCENT_TINT = RGBColor(0xF1, 0xDD, 0xD2)
MUTE        = RGBColor(0x78, 0x69, 0x5F)   # CANVAS 위 4.90:1
ON_ACCENT   = RGBColor(0xFF, 0xFF, 0xFF)
ON_ACCENT_S = RGBColor(0xF7, 0xE8, 0xDF)   # 어두워진 ACCENT 위 4.97:1
HAIRLINE    = RGBColor(0x9E, 0x8B, 0x79)   # CANVAS 위 3.04:1 — 장식선이라 4.5 미적용
BODY        = RGBColor(0x4D, 0x42, 0x3A)

# ── 타이포그래피 스케일 (16-great-ppt-design-research.md §2 + sol 검토 반영) ──
# 발표 투사 기준 본문 18pt 미만 금지(Microsoft 권고). 인접 단계 비율 약 1.3.
TYPE = {
    "display":  48,   # 표지 워드마크
    "title":    30,   # 콘텐츠 슬라이드 제목
    "subtitle": 18,   # 부제
    "label":    14,   # 챕터 라벨
    "body":     18,   # 본문·표 셀·불릿
    "caption":  12,   # 페이지 번호·보조 캡션
}

# ── 폰트 (냥BTI와 동일 — Srit이 한글 조합 보강한 빌드) ───────────────────────
F_LIGHT  = "Paperlogy Filled Light"
F_REG    = "Paperlogy Filled Regular"
F_MED    = "Paperlogy Filled Medium"
F_SEMI   = "Paperlogy Filled SemiBold"
F_BOLD   = "Paperlogy Filled Bold"
F_XBOLD  = "Paperlogy Filled ExtraBold"
F_BLACK  = "Paperlogy Filled Black"

FONT_DIR = pathlib.Path.home() / "Library" / "Fonts"
FONT_FILES = {
    F_LIGHT: FONT_DIR / "Paperlogy-Filled-3Light.ttf",
    F_REG:   FONT_DIR / "Paperlogy-Filled-4Regular.ttf",
    F_MED:   FONT_DIR / "Paperlogy-Filled-5Medium.ttf",
    F_SEMI:  FONT_DIR / "Paperlogy-Filled-6SemiBold.ttf",
    F_BOLD:  FONT_DIR / "Paperlogy-Filled-7Bold.ttf",
    F_XBOLD: FONT_DIR / "Paperlogy-Filled-8ExtraBold.ttf",
    F_BLACK: FONT_DIR / "Paperlogy-Filled-9Black.ttf",
}

# ── 마스터 앵커 (냥BTI 16:9 마스터 재사용) ───────────────────────────────────
SW, SH      = 10.0, 5.625
LX, W       = 0.50, 9.00
LABEL_Y     = 0.28
TITLE_Y     = 0.60
SUB_Y       = 1.15
CONTENT_Y   = 1.68
CONTENT_BOT = 5.08
FOOTER_Y    = 5.28

_violations = []


def tile(n, gap):
    """등폭 타일 규칙 — 절대 하드코딩하지 않는다."""
    return (W - (n - 1) * gap) / n


def _check(bottom, tag):
    if bottom > CONTENT_BOT + 1e-6:
        _violations.append(f"{tag}: bottom {bottom:.3f}\" > CONTENT_BOT {CONTENT_BOT}\"")


# ── measure_kr() — Pillow 사전 측정 (14-deck-tooling-research.md §4 adopt) ──
_PIL_FONT_CACHE = {}
_PX_PER_IN = 96  # 측정용 임의 해상도 — 상대 비율만 필요하므로 절대값은 무관


def _pil_font(family, size_pt):
    key = (family, round(size_pt, 1))
    if key in _PIL_FONT_CACHE:
        return _PIL_FONT_CACHE[key]
    path = FONT_FILES.get(family)
    if path is None or not path.exists():
        _PIL_FONT_CACHE[key] = None
        return None
    f = ImageFont.truetype(str(path), int(size_pt * _PX_PER_IN / 72))
    _PIL_FONT_CACHE[key] = f
    return f


_MEASURE_IMG = Image.new("RGB", (10, 10))
_MEASURE_DRAW = ImageDraw.Draw(_MEASURE_IMG)


def measure_kr(text, family, size_pt, width_in, line_spacing=1.18):
    """실제 렌더 전에 예상 줄바꿈 높이(in)를 계산한다.

    한국어는 띄어쓰기 기준 줄바꿈이 부정확해질 수 있어(조사·복합명사가 길게
    이어짐) 글자 단위로 감싼다 — 14-deck-tooling-research.md §1·§4가 확인한
    실제 패턴. 폰트를 못 찾으면 None을 반환하고, 호출자는 기존 고정값으로
    폴백한다(측정 실패가 빌드를 막지 않는다 — 이건 안전망이지 유일한 진실이
    아니다).
    """
    font = _pil_font(family, size_pt)
    if font is None:
        return None
    width_px = width_in * _PX_PER_IN
    lines, cur = [], ""
    for ch in text:
        trial = cur + ch
        w = _MEASURE_DRAW.textlength(trial, font=font)
        if w > width_px and cur:
            lines.append(cur)
            cur = ch
        else:
            cur = trial
    if cur:
        lines.append(cur)
    bbox = font.getbbox("가나다")
    line_h_px = (bbox[3] - bbox[1]) * line_spacing
    total_px = line_h_px * max(1, len(lines))
    return total_px / _PX_PER_IN


def font(run, family, size, color, spacing=None):
    """latin + ea + cs 모두 지정 — ea를 빼면 한글이 대체 폰트로 떨어진다."""
    f = run.font
    f.size = Pt(size)
    f.color.rgb = color
    f.name = family
    f.bold = False
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", family)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, family, size, color, first=False, space_before=0,
         line=None, align=None, spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    font(r, family, size, color, spacing)
    return p


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75, shape=MSO_SHAPE.RECTANGLE,
         radius=None, rotation=None):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    if rotation is not None:
        sh.rotation = rotation
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    return sh


def diamond(slide, cx, cy, d, fill):
    """15-deck-moodboard.md가 승인한 회전 다이아몬드 배지/도형."""
    return rect(slide, cx - d / 2, cy - d / 2, d, d, fill=fill,
                shape=MSO_SHAPE.DIAMOND)


def notes(slide, text):
    """발표자 노트 — sol 권고: 정정·출처 각주는 캔버스에서 글자를 줄이는
    대신 여기로 옮긴다. 청중이 보는 화면은 깨끗하게, 근거는 발표자만 본다."""
    slide.notes_slide.notes_text_frame.text = text


def hairline(slide, x, y, w, color=HAIRLINE, weight=0.9):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Emu(int(weight * 9525)))
    ln.shadow.inherit = False
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln


# ── 슬라이드 셸 ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
BLANK = prs.slide_layouts[6]
_page = {"n": 0}


def new_slide(kind="content", label=None, title=None, subtitle=None, numbered=True):
    s = prs.slides.add_slide(BLANK)
    _page["n"] += 1

    if kind == "content":
        rect(s, 0, 0, SW, SH, fill=CANVAS)
    elif kind == "divider":
        # 15가 승인한 구분 슬라이드 패턴 — ACCENT 단색 배경. 사진 배경 없음
        # (라이선스 사진 미확보, 15-deck-moodboard.md §차용하지 않을 것).
        rect(s, 0, 0, SW, SH, fill=ACCENT)

    on_dark = kind == "divider"
    if label:
        tf = textbox(s, LX, LABEL_Y, W, 0.26)
        para(tf, label, F_MED, TYPE["label"], ON_ACCENT_S if on_dark else MUTE,
             first=True, spacing=0.6)
    if title:
        tf = textbox(s, LX, TITLE_Y, W, 0.50)
        para(tf, title, F_BOLD, TYPE["title"], ON_ACCENT if on_dark else INK,
             first=True)
    if subtitle:
        tf = textbox(s, LX, SUB_Y, W, 0.36)
        para(tf, subtitle, F_LIGHT, TYPE["subtitle"],
             ON_ACCENT_S if on_dark else BODY, first=True)

    if numbered:
        tf = textbox(s, LX, FOOTER_Y, 1.0, 0.22)
        para(tf, f"{_page['n']:02d}", F_MED, TYPE["caption"],
             ON_ACCENT_S if on_dark else MUTE, first=True)
    return s


def bullets(slide, items, x=LX, y=CONTENT_Y, w=W, size=None, color=INK,
            family=F_REG, gap=0.36, tag="bullets"):
    size = size or TYPE["body"]
    cur = y
    for it in items:
        indent = it.get("indent", 0)
        fam = {"b": F_SEMI, "r": F_REG, "l": F_LIGHT}[it.get("w", "r")]
        col = it.get("color", color)
        sz = it.get("size", size)
        marker = "·" if indent else "—"
        tf = textbox(slide, x + 0.14 * indent, cur, w - 0.14 * indent, 0.32)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"{marker}  "
        font(r, fam, sz, ACCENT if not indent else MUTE)
        r2 = p.add_run(); r2.text = it["t"]
        font(r2, fam, sz, col)
        cur += it.get("gap", gap)
    _check(cur, tag)
    return cur


def _table_block_height(cols, rows, widths, size, hdr_size, row_h):
    """table()과 같은 계산을 실제로 그리지 않고 먼저 구해 — 세로 중앙
    배치(vcenter)에 쓴다. 16-great-ppt-design-research.md §1·§6 권고:
    "행 높이를 빈 공간에 맞춰 늘리지 말고, 표 그룹의 위치만 옮겨라."."""
    total = 0.40  # 헤더 행 + 헤더 밑줄 여백, table()의 y+0.40과 동일
    for row in rows:
        needed = row_h
        for i, cell in enumerate(row):
            m = measure_kr(cell, F_REG, size, widths[i] - 0.10)
            if m is not None:
                needed = max(needed, m + 0.20)
        total += needed
    return total


def table(slide, cols, rows, y=CONTENT_Y, widths=None, size=None, hdr_size=None,
          row_h=0.62, tag="table", row_anchor=MSO_ANCHOR.TOP,
          vcenter_in=None):
    """하드코딩 row_h 대신 measure_kr()로 각 행의 실제 필요 높이를 먼저 계산한다
    (14-deck-tooling-research.md §4 권고 — 고정 상한 검사 하나보다 한 단계
    앞선 preflight). 측정 실패 시(폰트 못 찾음) row_h로 폴백.

    vcenter_in=(top, bottom)을 주면 표 그룹 전체를 그 구간 안에서 세로
    중앙에 놓는다 — 짧은 표가 콘텐츠 영역 위쪽에만 몰려 아래가 빈 것처럼
    보이는 문제의 실제 원인은 행 높이가 아니라 표 그룹의 y 위치였다
    (sol 검토, 2026-08-06)."""
    size = size or TYPE["body"]
    hdr_size = hdr_size or TYPE["body"]
    n = len(cols)
    widths = widths or [W / n] * n
    if vcenter_in is not None:
        top, bottom = vcenter_in
        block_h = _table_block_height(cols, rows, widths, size, hdr_size, row_h)
        y = top + max(0, (bottom - top - block_h) / 2)
    xs, acc = [], LX
    for cw in widths:
        xs.append(acc); acc += cw

    for i, c in enumerate(cols):
        tf = textbox(slide, xs[i], y, widths[i] - 0.10, 0.28)
        para(tf, c, F_SEMI, hdr_size, ACCENT, first=True)
    hairline(slide, LX, y + 0.32, W, color=HAIRLINE, weight=1.2)

    cur = y + 0.40
    for r_i, row in enumerate(rows):
        needed = row_h
        for i, cell in enumerate(row):
            m = measure_kr(cell, F_REG, size, widths[i] - 0.10)
            if m is not None:
                needed = max(needed, m + 0.20)
        h = needed
        for i, cell in enumerate(row):
            fam = F_SEMI if (i == 0) else F_REG
            col = INK if (i == 0) else BODY
            tf = textbox(slide, xs[i], cur, widths[i] - 0.10, h,
                         anchor=row_anchor)
            para(tf, cell, fam, size, col, first=True, line=1.18)
        cur += h
        if r_i < len(rows) - 1:
            hairline(slide, LX, cur - 0.06, W)
    _check(cur, tag)
    return cur


OUT = pathlib.Path(__file__).parent / "music-diary-presentation.pptx"

# ══════════════════════════════════════════════════════════════════════════
# 01 · 표지 — CANVAS 배경(정정: 이전 버전은 ACCENT 배경으로 잘못 표기됐었음,
#      그건 구분 슬라이드 패턴이다 — presentation-outline.md 01번 정정 참고)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", numbered=False)
tf = textbox(s, LX, 0.42, W, 0.26)
para(tf, "AI 활용 웹콘텐츠 실무역량 과정 · 팀 Clova", F_SEMI, 11, MUTE, first=True,
     spacing=1.2)
tf = textbox(s, LX, 1.72, 7.4, 1.0)
para(tf, "Music Diary", F_BLACK, 52, INK, first=True)
tf = textbox(s, LX, 2.86, 8.2, 0.34)
para(tf, "음악 일기 — 오늘의 기분으로 내 음악을 다시 만나는 로컬 플레이어",
     F_LIGHT, 15, BODY, first=True)
hairline(s, LX, 3.46, 2.2, color=ACCENT_TINT, weight=1.4)
tf = textbox(s, LX, 3.72, 4.0, 0.80)
para(tf, "2026-08-14", F_LIGHT, 15, MUTE, first=True, line=1.6)
para(tf, "발표팀  Clova", F_LIGHT, 15, MUTE, line=1.6)
# 사진 콜라주 없음 — 도형 클러스터로 표지 리듬만 재현 (15 §차용하지 않을 것)
diamond(s, 8.65, 4.05, 1.35, ACCENT)
diamond(s, 9.55, 3.30, 0.85, ACCENT_TINT)

# ══════════════════════════════════════════════════════════════════════════
# 02 · 문제와 주 페르소나 (8분 1:00 / 5분 0:45)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "01 · 문제",
              "“그 곡을 왜 반복해서 들었는지, 나중엔 기억이 안 나요”",
              "이서준 · 19세 · 고등학교 3학년 — 주 페르소나")
table(s, ["지금", "Music Diary가 채우는 것"],
      [["플레이리스트를 기분별로 나눠 관리하지만 \"왜 그날 이 곡을 반복했는지\"는 시간이 지나면 잊힘",
        "기분 선택 → 재생 → 재생 직후 한 줄 일기로 그날의 맥락이 남음"],
       ["스트리밍 앱의 \"오늘의 추천\"은 자신의 기분과 무관하게 인기곡 위주",
        "추천은 내 라이브러리 안에서만, 기분 태그로 직접 필터링"]],
      y=CONTENT_Y + 0.08, widths=[4.75, 4.25], row_h=1.20,
      row_anchor=MSO_ANCHOR.MIDDLE, tag="02 table")

# ══════════════════════════════════════════════════════════════════════════
# 03 · 조사에서 정한 제품 위치 (8분 0:45 / 5분 0:10·한 문장)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "02 · 포지셔닝", "카탈로그 경쟁이 아니라, 내 파일 + 그날의 기분",
              "Spotify·YouTube Music·Melon 등 5개 서비스 조사 후 정한 차별점")
table(s, ["비교 축", "스트리밍 서비스 5종", "Music Diary"],
      [["핵심 경쟁 축",
        "카탈로그·발견·구독 경험 (5종이 저마다 다른 방식으로 경쟁)",
        "내 파일의 감정 맥락"],
       ["라이브러리", "회사의 카탈로그", "사용자 자신의 파일"],
       ["감정 기록", "없음 (5종 전부 확인)", "재생 직후 한 줄 일기"]],
      y=CONTENT_Y + 0.08, widths=[1.55, 3.85, 3.60], row_h=0.95,
      row_anchor=MSO_ANCHOR.MIDDLE, tag="03 table")
notes(s, "\"대규모 청취 데이터 기반 추천\"은 5종 공통이 아니라 Spotify 한정 — "
         "다른 4종의 차별점은 큐 예측가능성(YouTube Music)·차트/팬덤(Melon)·"
         "음질/정리(Apple Music)·구독 복잡도(Amazon Music)로 각기 다름.")

# ══════════════════════════════════════════════════════════════════════════
# 04 · 해결 흐름과 범위 (8분 0:35 / 5분 0:25) — 새 제안, 15에 없음, 팀 승인 필요
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "새 제안 · 15 미승인", "기분 선택 → 재생 → 한 줄 일기 → 기록",
              "MVP 5탭(홈·라이브러리·탐색·일기·마이)이 담는 사용자 경로 하나")
gap = 0.24
bw = tile(4, gap)
steps = [("01", "기분 선택", "오늘의 기분\n5종 중 하나"),
         ("02", "재생", "라이브러리 안에서\n필터링된 곡 재생"),
         ("03", "한 줄 일기", "재생 종료 직후\n그날 기분 기록"),
         ("04", "날짜별 기록", "히스토리로 누적\n나중에 돌아봄")]
BOX_Y, BOX_H = 2.30, 1.55
for i, (num, lab, desc) in enumerate(steps):
    x = LX + i * (bw + gap)
    box = rect(s, x, BOX_Y, bw, BOX_H, fill=ACCENT)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, num, F_XBOLD, 18, ON_ACCENT_S, first=True)
    para(tf, lab, F_SEMI, 18, ON_ACCENT, space_before=4)
    para(tf, desc, F_LIGHT, 15, ON_ACCENT_S, space_before=4, line=1.25)
    if i < len(steps) - 1:
        ay = BOX_Y + BOX_H / 2
        rect(s, x + bw + 0.04, ay - 0.02, gap - 0.08, 0.04, fill=ACCENT)
_check(BOX_Y + BOX_H, "04 flow")

# ══════════════════════════════════════════════════════════════════════════
# 05 · AI 사용과 팀 검증 (8분 0:50 / 5분 0:35)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "03 · AI 프로세스", "AI 결과를 그대로 쓰지 않았습니다",
              "조사 모델과 검토 모델을 분리 — 같은 모델이 자기 작업을 검토하면 결함을 못 찾는다는 전제")
table(s, ["AI가 한 것", "검토가 잡아낸 것"],
      [["`07`·`08` 계획을 검토 모델(sol)에게 사실 검증시킴",
        "\"일정 안에 못 만든다\"는 답 — 그대로 기록(BLOCKER 4건 포함 18건, `01`)"],
       ["컴포넌트 패턴 조사(luna, `02`)",
        "다른 모델(sol)이 재검증 — 깨진 링크·labs(실험적) 출처 오류를 잡음(`03`)"],
       ["번들 오디오 라이선스 후보 조사",
        "Pixabay가 CC가 아니라 자체 라이선스·재배포 금지임을 발견해 제외"]],
      y=CONTENT_Y + 0.08, widths=[3.75, 5.25], row_h=0.95,
      row_anchor=MSO_ANCHOR.MIDDLE, tag="05 table")
notes(s, "이퀄라이저·크로스페이드·태그 자동분류는 취향이 아니라 각각 Web Audio "
         "그래프·두 번째 오디오 요소·ID3 파서가 필요해서 현재 구조에서 뺀 것입니다.")

# ══════════════════════════════════════════════════════════════════════════
# 06 · 라이브 데모 (8분 2:00 / 5분 1:30) — 15가 승인한 구분 슬라이드 패턴
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("divider", numbered=True)
# 정정(sol 검토) — 원래 다이아몬드(중앙, d=1.9")가 제목·부제 텍스트 박스와
# 겹쳐 대비가 무너졌다. 우상단 모서리로 옮겨 텍스트 안전 영역 밖에 둠.
diamond(s, 8.25, 1.35, 1.20, ACCENT_TINT)
tf = textbox(s, LX + 0.30, 2.10, W - 0.60, 0.90, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
para(tf, "[코딩 및 통합 검증 후 확정]", F_BLACK, 26, ON_ACCENT, first=True,
     align=PP_ALIGN.CENTER, line=1.2)
tf = textbox(s, LX + 0.30, 3.30, W - 0.60, 0.40, anchor=MSO_ANCHOR.TOP,
             align=PP_ALIGN.CENTER)
para(tf, "실제 화면과 90~150초 시나리오가 리허설로 확인되기 전까지 이 자리 유지",
     F_LIGHT, TYPE["caption"], ON_ACCENT_S, first=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# 07 · 실제 차이와 정직한 한계 (8분 1:10, 08과 합쳐 5분 0:45)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "04 · 한계", "정직하게 말씀드리는 제약",
              "계정·백엔드가 없다는 것은 제약이자 약속")
table(s, ["목표", "감수하는 것 — 이번 MVP의 결정"],
      [["목표: 파일을 서버에 전송하지 않는 로컬 처리 (구현 후 네트워크 탭으로 검증 예정)",
        "새로고침 후 파일을 다시 선택해야 함 — 파일 본문 대신 메타데이터만 저장하기로 한 결정"],
       ["계정 없이 바로 사용", "다른 기기로 옮기면 기록이 이어지지 않음"],
       ["라이브러리는 사용자 자신의 파일",
        "장르·앨범 자동 분류 없음 — ID3/MP4 태그 파서를 범위에 넣지 않은 결정"]],
      y=CONTENT_Y + 0.08, widths=[3.75, 5.25], row_h=0.84,
      row_anchor=MSO_ANCHOR.MIDDLE, tag="07 table")

# ══════════════════════════════════════════════════════════════════════════
# 08 · 3인 역할 · 다음 검증 · 마무리 (8분 0:45, 5분은 07과 통합)
# ══════════════════════════════════════════════════════════════════════════
s = new_slide("content", "05 · 팀과 다음 단계", "[팀 확정 필요 — 3인 역할 분담]",
              "다음 검증: 반 친구 최소 3명 · 5분 프로토타입 테스트")
# 팀 사진 자리 — 15가 승인한 원형 자리 패턴, 실제 사진은 팀이 준비.
# 정정(sol 검토) — 원래 왼쪽에 몰려 오른쪽 절반이 비어 보였다. 콘텐츠 폭
# 9.00" 전체에 3등분으로 재배치(순수 레이아웃 수정 — placeholder 문구
# 자체는 팀 역할·사진이 아직 미정이라 그대로 둠, 여기서 지어내지 않음).
PHOTO_Y, PHOTO_D = CONTENT_Y + 0.06, 1.10
for i in range(3):
    cx = LX + W * (i + 0.5) / 3
    circ = rect(s, cx - PHOTO_D / 2, PHOTO_Y, PHOTO_D, PHOTO_D, fill=ACCENT_TINT,
                shape=MSO_SHAPE.OVAL)
    tfc = circ.text_frame
    tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tfc, "[팀원 사진]", F_MED, 12, ACCENT, first=True, align=PP_ALIGN.CENTER)
cur = PHOTO_Y + PHOTO_D + 0.24
end = bullets(s, [
    {"t": "역할 분담 — [미정, 팀에서 채울 것]", "w": "b"},
    {"t": "다음 검증 1 — 기분 선택 → 추천곡 흐름이 직관적인지", "indent": 1},
    {"t": "다음 검증 2 — 일기를 실제로 남길 동기가 생기는지", "indent": 1},
    {"t": "다음 검증 3 — 파일을 새로고침마다 다시 선택하는 게 얼마나 불편한지", "indent": 1},
], y=cur, gap=0.34, tag="08 bullets")
tf = textbox(s, LX, end + 0.12, W, 0.40)
para(tf, "“오늘의 기분을 고르면, 내 음악이 그날을 기억해 줍니다.”",
     F_SEMI, 20, ACCENT, first=True)
_check(end + 0.52, "08 closing")

# ── 저장 + 검증 ────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"saved: {OUT}  ({OUT.stat().st_size:,} bytes)")
print(f"slides: {len(prs.slides._sldIdLst)}")

# 저장 후 실제 도형을 전부 훑어 재검증 — 냥BTI와 같은 사후 안전장치.
EMU_IN = 914400
hard = []
for idx, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        if sh.top is None or sh.width is None or sh.height is None:
            continue
        top, bot, wid = sh.top / EMU_IN, (sh.top + sh.height) / EMU_IN, sh.width / EMU_IN
        if wid > 9.9:                      # 전면 배경
            continue
        if abs(top - FOOTER_Y) < 0.02:     # 푸터 행 — 규격상 정상
            continue
        if bot > CONTENT_BOT + 1e-3:
            txt = sh.text_frame.text[:30].replace("\n", "/") if sh.has_text_frame else ""
            hard.append(f"slide {idx}: bottom {bot:.3f}\" — {txt}")

if _violations or hard:
    print(f"\n!!! CONTENT_BOT violations — build-time {len(_violations)}, post-save {len(hard)}")
    for v in _violations + hard:
        print("   ", v)
    raise SystemExit(1)
print(f"CONTENT_BOT {CONTENT_BOT}\" — verified across every shape on all "
      f"{len(prs.slides._sldIdLst)} slides")
