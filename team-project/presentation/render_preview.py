#!/usr/bin/env python3
"""build_deck.py 결과물을 실제로 눈으로 보기 위한 근사 렌더러.

이 기기에 PowerPoint·LibreOffice가 없어 진짜 렌더가 불가능하다(README 확인).
python-pptx가 저장한 .pptx를 다시 읽어 도형 위치·색·텍스트를 그대로 좌표
변환해 PIL로 그린다 — PowerPoint의 정확한 타이포그래피 렌더링(커닝, 실제
줄바꿈 엔진)까지 재현하지는 않지만, 레이아웃·색·여백·정렬을 검토하기에는
충분한 근사치다. 리뷰·QA 전용 스크립트 — 산출물(.pptx)에는 영향 없음.

실행:  <venv>/bin/python render_preview.py
출력:  preview/slide-01.png ... slide-08.png
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont
import pathlib

HERE = pathlib.Path(__file__).parent
PPTX = HERE / "music-diary-presentation.pptx"
OUT_DIR = HERE / "preview"
OUT_DIR.mkdir(exist_ok=True)

SCALE = 160  # px per inch — 10"x5.625" 마스터 → 1600x900
FONT_DIR = pathlib.Path.home() / "Library" / "Fonts"
FONT_MAP = {
    "Paperlogy Filled Light":     FONT_DIR / "Paperlogy-Filled-3Light.ttf",
    "Paperlogy Filled Regular":   FONT_DIR / "Paperlogy-Filled-4Regular.ttf",
    "Paperlogy Filled Medium":    FONT_DIR / "Paperlogy-Filled-5Medium.ttf",
    "Paperlogy Filled SemiBold":  FONT_DIR / "Paperlogy-Filled-6SemiBold.ttf",
    "Paperlogy Filled Bold":      FONT_DIR / "Paperlogy-Filled-7Bold.ttf",
    "Paperlogy Filled ExtraBold": FONT_DIR / "Paperlogy-Filled-8ExtraBold.ttf",
    "Paperlogy Filled Black":     FONT_DIR / "Paperlogy-Filled-9Black.ttf",
}
_font_cache = {}


def get_font(name, size_pt):
    key = (name, round(size_pt))
    if key in _font_cache:
        return _font_cache[key]
    path = FONT_MAP.get(name)
    px = max(1, int(size_pt * SCALE / 72))
    f = ImageFont.truetype(str(path), px) if path and path.exists() else ImageFont.load_default()
    _font_cache[key] = f
    return f


def emu_to_px(emu):
    return emu / 914400 * SCALE


def rgb(color_format):
    try:
        c = color_format.rgb
        return (c[0], c[1], c[2])
    except Exception:
        return None


prs = Presentation(PPTX)
for idx, slide in enumerate(prs.slides, 1):
    W, H = int(emu_to_px(prs.slide_width)), int(emu_to_px(prs.slide_height))
    img = Image.new("RGB", (W, H), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        if shape.left is None:
            continue
        x, y = emu_to_px(shape.left), emu_to_px(shape.top)
        w, h = emu_to_px(shape.width), emu_to_px(shape.height)

        # 채우기
        fill_rgb = None
        try:
            if shape.fill.type is not None:
                fill_rgb = rgb(shape.fill.fore_color)
        except Exception:
            pass
        ast = None
        try:
            ast = shape.auto_shape_type
        except Exception:
            pass
        if fill_rgb:
            if ast is not None and str(ast).startswith("OVAL"):
                draw.ellipse([x, y, x + w, y + h], fill=fill_rgb)
            elif ast is not None and "DIAMOND" in str(ast):
                cx, cy = x + w / 2, y + h / 2
                draw.polygon([(cx, y), (x + w, cy), (cx, y + h), (x, cy)], fill=fill_rgb)
            else:
                draw.rectangle([x, y, x + w, y + h], fill=fill_rgb)

        # 텍스트 — 먼저 전체 줄바꿈과 높이를 계산한 다음, vertical_anchor
        # (TOP/MIDDLE/BOTTOM)에 맞춰 시작 y를 정한다. python-pptx의
        # MSO_ANCHOR를 무시하고 항상 위에서부터 그리면, 실제로는 세로 중앙
        # 정렬된 표 셀(table()의 row_anchor=MIDDLE)이 이 근사 렌더에서만
        # 위쪽에 몰려 보이는 착시가 생긴다 — 실제 .pptx 데이터는 정확해도
        # 리뷰용 이미지가 잘못된 인상을 준다.
        if shape.has_text_frame:
            para_lines = []  # [(size, color, fam, [line, ...], align), ...]
            block_h = 0
            for para in shape.text_frame.paragraphs:
                if not para.runs:
                    para_lines.append((12, (30, 30, 30), None, [""], para.alignment))
                    block_h += 14
                    continue
                r = para.runs[0]
                fam = r.font.name or "Paperlogy Filled Regular"
                size = r.font.size.pt if r.font.size else 12
                color = rgb(r.font.color) if r.font.color and r.font.color.type else (30, 30, 30)
                text = "".join(run.text for run in para.runs)
                f = get_font(fam, size)
                # 아주 근사한 글자 단위 줄바꿈 (실제 word_wrap 근사).
                # 런 텍스트 자체에 명시적 개행이 섞여 있을 수 있어 먼저 분리한다
                # (Pillow의 textlength는 멀티라인 문자열을 받으면 예외를 던진다).
                lines = []
                for hard_line in text.split("\n"):
                    cur = ""
                    for ch in hard_line:
                        trial = cur + ch
                        tw = draw.textlength(trial, font=f)
                        if tw > w and cur:
                            lines.append(cur)
                            cur = ch
                        else:
                            cur = trial
                    lines.append(cur)
                para_lines.append((size, color, fam, lines, para.alignment))
                block_h += len(lines) * size * SCALE / 72 * 1.22 + 2

            anchor = shape.text_frame.vertical_anchor
            if anchor == MSO_ANCHOR.MIDDLE:
                ty = y + max(0, (h - block_h) / 2)
            elif anchor == MSO_ANCHOR.BOTTOM:
                ty = y + max(0, h - block_h)
            else:
                ty = y

            for size, color, fam, lines, align in para_lines:
                f = get_font(fam, size) if fam else get_font("Paperlogy Filled Regular", size)
                for ln in lines:
                    lw = draw.textlength(ln, font=f)
                    lx = x
                    if align == PP_ALIGN.CENTER:
                        lx = x + (w - lw) / 2
                    elif align == PP_ALIGN.RIGHT:
                        lx = x + w - lw
                    draw.text((lx, ty), ln, font=f, fill=color)
                    ty += size * SCALE / 72 * 1.22
                ty += 2

    out = OUT_DIR / f"slide-{idx:02d}.png"
    img.save(out)
    print(f"wrote {out}")
